import io
from collections import OrderedDict
from copy import deepcopy
from difflib import SequenceMatcher
from itertools import combinations
from pathlib import Path
import re
from typing import Dict, List, Optional, Tuple

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt

APP_DIR = Path(__file__).resolve().parent

REPORT_TYPES: Dict[str, Dict[str, object]] = {
    "Lam Research": {
        "id": "lam",
        "template_candidates": [APP_DIR / "template.pptx"],
    },
    "\ud55c\uad6d\uc54c\ucf5c\uadf8\ub8f9": {
        "id": "kai",
        "template_candidates": [
            APP_DIR / "kai-template.pptx",
            APP_DIR.parent / "KAI-auto" / "kai-template.pptx",
            APP_DIR.parent / "KAI-auto" / "template.pptx",
        ],
    },
}


def shape_index_from_name(name: str) -> Optional[int]:
    match = re.search(r"(\d+)$", name.strip())
    return int(match.group(1)) if match else None


def resolve_template_path(report_label: str) -> Path:
    report_config = REPORT_TYPES.get(report_label)
    if report_config is None:
        raise ValueError(f"지원하지 않는 보고서 유형입니다: {report_label}")

    candidates = [Path(path) for path in report_config.get("template_candidates", [])]
    for candidate in candidates:
        if candidate.exists():
            return candidate

    checked = ", ".join(str(path) for path in candidates)
    raise FileNotFoundError(f"템플릿 파일을 찾지 못했습니다. 확인 경로: {checked}")


def identify_question_columns(df: pd.DataFrame, required: int = 10) -> List[str]:
    normalized = {str(col).strip().lower(): col for col in df.columns}
    selected: List[str] = []

    for idx in range(1, required + 1):
        candidates = {f"q{idx}", f"q{idx}.", f"q{idx} ", f"문항{idx}", f"문항 {idx}", f"{idx}번", f"{idx}"}
        found = None
        for key, original in normalized.items():
            if key in candidates or re.fullmatch(rf"q\s*{idx}", key):
                found = original
                break
        if found and found not in selected:
            selected.append(found)

    if len(selected) >= required:
        return selected[:required]

    score_like_cols: List[str] = []
    for col in df.columns:
        series = pd.to_numeric(df[col], errors="coerce").dropna()
        if not series.empty and series.between(1, 4).all():
            score_like_cols.append(col)

    for col in score_like_cols:
        if col not in selected:
            selected.append(col)
        if len(selected) == required:
            break

    if len(selected) < required:
        raise ValueError("1~4점 척도 문항 컬럼 10개를 찾지 못했습니다. 엑셀 컬럼명을 확인해 주세요.")

    return selected


def to_100_scale(raw_avg: float) -> float:
    return (raw_avg / 4.0) * 100.0


def format_score(value: float) -> str:
    return f"{value:.1f}"


def compute_lam_metrics(
    df: pd.DataFrame, question_cols: List[str]
) -> Tuple[Dict[str, float], Dict[int, Dict[int, float]], Dict[int, Dict[int, int]], int]:
    numeric = df[question_cols].apply(pd.to_numeric, errors="coerce")
    respondent_count = int(numeric.dropna(how="all").shape[0])

    sub_avgs_raw = [numeric[col].mean() for col in question_cols]
    sub_avgs_100 = [to_100_scale(v) if pd.notna(v) else 0.0 for v in sub_avgs_raw]

    total_avg_00 = to_100_scale(numeric.mean(axis=1).mean()) if not numeric.empty else 0.0
    total_avg_01 = to_100_scale(numeric.iloc[:, 0:5].mean(axis=1).mean())
    total_avg_02 = to_100_scale(numeric.iloc[:, 5:8].mean(axis=1).mean())
    total_avg_03 = to_100_scale(numeric.iloc[:, 8:10].mean(axis=1).mean())
    total_avg_04 = to_100_scale(numeric.iloc[:, 1:5].mean(axis=1).mean())

    placeholders: Dict[str, float] = {
        "total_avg_00": total_avg_00,
        "total_avg_01": total_avg_01,
        "total_avg_02": total_avg_02,
        "total_avg_03": total_avg_03,
        "total_avg_04": total_avg_04,
    }

    for idx, value in enumerate(sub_avgs_100, start=1):
        placeholders[f"sub_avg_{idx:02d}"] = value

    percentages_by_question: Dict[int, Dict[int, float]] = {}
    counts_by_question: Dict[int, Dict[int, int]] = {}

    for q_idx, col in enumerate(question_cols, start=1):
        answers = pd.to_numeric(numeric[col], errors="coerce").dropna().astype(int)
        total = len(answers)
        counts = {score: int((answers == score).sum()) for score in [1, 2, 3, 4]}
        percentages = {score: ((counts[score] / total) * 100.0 if total else 0.0) for score in [1, 2, 3, 4]}
        counts_by_question[q_idx] = counts
        percentages_by_question[q_idx] = percentages

    return placeholders, percentages_by_question, counts_by_question, respondent_count


def set_text_preserve_style(text_frame, value: str) -> None:
    if not text_frame.paragraphs:
        text_frame.text = value
        return

    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = value
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = value

    for extra_paragraph in text_frame.paragraphs[1:]:
        for run in extra_paragraph.runs:
            run.text = ""


def replace_text_placeholders(prs: Presentation, replacements: Dict[str, str]) -> None:
    def replace_in_runs(paragraph) -> None:
        if not paragraph.runs:
            return
        for run in paragraph.runs:
            updated = run.text
            for key, value in replacements.items():
                updated = updated.replace(key, value)
            if updated != run.text:
                run.text = updated

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    replace_in_runs(paragraph)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            replace_in_runs(paragraph)


def set_chart_number_format(chart, number_format: str) -> None:
    for plot in chart.plots:
        labels = plot.data_labels
        labels.show_value = True
        labels.number_format = number_format


def update_lam_chart_0(shape, placeholders: Dict[str, float]) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = [
        "전체 평균",
        "과정만족도 평균",
        "1",
        "2",
        "3",
        "4",
        "5",
        "강사만족도 평균",
        "6",
        "7",
        "8",
        "운영 만족도 평균",
        "9",
        "10",
    ]
    chart_data.add_series(
        "계열",
        (
            round(placeholders["total_avg_00"], 1),
            round(placeholders["total_avg_01"], 1),
            round(placeholders["sub_avg_01"], 1),
            round(placeholders["sub_avg_02"], 1),
            round(placeholders["sub_avg_03"], 1),
            round(placeholders["sub_avg_04"], 1),
            round(placeholders["sub_avg_05"], 1),
            round(placeholders["total_avg_02"], 1),
            round(placeholders["sub_avg_06"], 1),
            round(placeholders["sub_avg_07"], 1),
            round(placeholders["sub_avg_08"], 1),
            round(placeholders["total_avg_03"], 1),
            round(placeholders["sub_avg_09"], 1),
            round(placeholders["sub_avg_10"], 1),
        ),
    )
    shape.chart.replace_data(chart_data)
    set_chart_number_format(shape.chart, "0.0")


def update_lam_question_chart(shape, question_idx: int, percentages_by_question: Dict[int, Dict[int, float]]) -> None:
    percentages = percentages_by_question[question_idx]
    chart_data = CategoryChartData()
    chart_data.categories = ["100점", "75점", "50점", "25점"]
    chart_data.add_series(
        "계열",
        (
            percentages[4] / 100.0,
            percentages[3] / 100.0,
            percentages[2] / 100.0,
            percentages[1] / 100.0,
        ),
    )
    shape.chart.replace_data(chart_data)
    set_chart_number_format(shape.chart, "0%")


def set_font_cjk(run, font_name: str, font_size_pt: float) -> None:
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)


def format_table_font(table, font_name: str, font_size_pt: float) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_font_cjk(run, font_name, font_size_pt)


def update_lam_question_table(shape, question_idx: int, counts_by_question: Dict[int, Dict[int, int]]) -> None:
    table = shape.table
    counts = counts_by_question[question_idx]
    total = sum(counts.values())

    def pct(score: int) -> int:
        return int(round((counts[score] / total) * 100)) if total else 0

    rows = [
        f"{counts[4]}명({pct(4)}%)",
        f"{counts[3]}명({pct(3)}%)",
        f"{counts[2]}명({pct(2)}%)",
        f"{counts[1]}명({pct(1)}%)",
        f"{total}명(100%)",
    ]

    for offset, value in enumerate(rows, start=1):
        if len(table.rows) > offset and len(table.columns) > 1:
            set_text_preserve_style(table.cell(offset, 1).text_frame, value)


def summarize_qualitative_responses(series: pd.Series) -> str:
    deduped_counts: "OrderedDict[str, int]" = OrderedDict()
    filtered_count = 0

    def normalize_for_compare(value: str) -> str:
        lowered = re.sub(r"\s+", " ", value).strip().lower()
        return re.sub(r"[^0-9a-z가-힣\s]", "", lowered)

    def extract_keywords(value: str) -> set:
        tokens = re.findall(r"[0-9a-z가-힣]+", value)
        return {token for token in tokens if len(token) >= 2}

    def is_filtered(value: str) -> bool:
        normalized = re.sub(r"\s+", "", value).lower().strip().strip("\"'")
        normalized_wo_punct = normalized.rstrip(".!~")
        return normalized in {"", "."} or normalized_wo_punct in {
            "없음", "없습니다", "없다", "없어요", "none", "na", "n/a", "x", "ㄴ", "-"
        }

    def char_bigrams(value: str) -> set:
        compact = re.sub(r"\s+", "", value)
        if len(compact) < 2:
            return {compact}
        return {compact[idx : idx + 2] for idx in range(len(compact) - 1)}

    for raw in series.tolist():
        if pd.isna(raw):
            continue

        value = str(raw).strip()
        if is_filtered(value):
            filtered_count += 1
            continue

        key = re.sub(r"\s+", " ", value)
        deduped_counts[key] = deduped_counts.get(key, 0) + 1

    if not deduped_counts:
        return f"(유효 응답 없음)\n필터링 응답 {filtered_count}개"
    items = [
        {
            "text": text,
            "count": count,
            "norm": normalize_for_compare(text),
            "grams": char_bigrams(text),
            "keywords": extract_keywords(normalize_for_compare(text)),
        }
        for text, count in deduped_counts.items()
    ]

    parent = list(range(len(items)))

    def find(x: int) -> int:
        while parent[x] != x:
            parent[x] = parent[parent[x]]
            x = parent[x]
        return x

    def union(a: int, b: int) -> None:
        root_a = find(a)
        root_b = find(b)
        if root_a != root_b:
            parent[root_b] = root_a

    for left, right in combinations(range(len(items)), 2):
        left_item = items[left]
        right_item = items[right]

        if left_item["norm"] in right_item["norm"] or right_item["norm"] in left_item["norm"]:
            union(left, right)
            continue

        left_keywords = left_item["keywords"]
        right_keywords = right_item["keywords"]
        shared_keywords = left_keywords & right_keywords
        keyword_union = left_keywords | right_keywords
        keyword_similarity = len(shared_keywords) / len(keyword_union) if keyword_union else 0.0

        if shared_keywords and (
            any(len(keyword) >= 3 for keyword in shared_keywords) or keyword_similarity >= 0.5
        ):
            union(left, right)
            continue

        union_size = len(left_item["grams"] | right_item["grams"])
        bigram_similarity = (
            len(left_item["grams"] & right_item["grams"]) / union_size if union_size else 0.0
        )
        sequence_similarity = SequenceMatcher(None, left_item["norm"], right_item["norm"]).ratio()
        if max(bigram_similarity, sequence_similarity) >= 0.55 or (
            shared_keywords and max(bigram_similarity, sequence_similarity) >= 0.45
        ):
            union(left, right)

    grouped: "OrderedDict[int, List[Dict[str, object]]]" = OrderedDict()
    for idx, item in enumerate(items):
        root = find(idx)
        grouped.setdefault(root, []).append(item)

    cluster_summaries = []
    misc_items: List[Dict[str, object]] = []
    misc_total = 0

    theme_no = 1
    for group_items in grouped.values():
        group_total = sum(int(entry["count"]) for entry in group_items)
        representative = sorted(
            group_items,
            key=lambda entry: (int(entry["count"]), len(str(entry["text"]))),
            reverse=True,
        )[0]
        details = sorted(
            group_items,
            key=lambda entry: (int(entry["count"]), len(str(entry["text"]))),
            reverse=True,
        )

        if group_total == 1:
            misc_total += group_total
            misc_items.extend(details)
            continue

        cluster_summaries.append(
            {
                "representative": representative,
                "group_total": group_total,
                "details": details,
            }
        )

    cluster_summaries = sorted(
        cluster_summaries,
        key=lambda summary: (
            int(summary["group_total"]),
            len(str(summary["representative"]["text"])),
        ),
        reverse=True,
    )

    lines: List[str] = []
    for summary in cluster_summaries:
        representative = summary["representative"]
        group_total = int(summary["group_total"])
        details = summary["details"]

        lines.append(f"주제{theme_no}. {representative['text']} ({group_total})")
        for detail in details:
            text = str(detail["text"])
            count = int(detail["count"])
            lines.append(f"- {text} ({count})" if count > 1 else f"- {text}")
        theme_no += 1

    if misc_items:
        lines.append(f"기타 ({misc_total})")
        for detail in sorted(misc_items, key=lambda entry: str(entry["text"])):
            lines.append(f"- {detail['text']}")

    lines.append("")
    lines.append(f"필터링 응답 {filtered_count}개")
    return "\n".join(lines)


def update_lam_qualitative_table(shape, df: pd.DataFrame) -> None:
    table = shape.table
    if len(df.columns) < 14:
        raise ValueError("L~N열(12~14번째 열)을 찾지 못했습니다. 엑셀 컬럼 구성을 확인해 주세요.")

    qualitative_cols = [df.iloc[:, 11], df.iloc[:, 12], df.iloc[:, 13]]

    for idx, column in enumerate(qualitative_cols, start=1):
        summary_text = summarize_qualitative_responses(column)
        if len(table.rows) > 1 and len(table.columns) > idx:
            set_text_preserve_style(table.cell(1, idx).text_frame, summary_text)


def parse_kai_excel(excel_bytes: bytes) -> Tuple[pd.DataFrame, List[int], Optional[int], List[int]]:
    workbook = pd.read_excel(io.BytesIO(excel_bytes), sheet_name=None, header=None)
    if not workbook:
        raise ValueError("엑셀 파일이 비어 있습니다.")

    target_sheet_name = next(
        (name for name in workbook.keys() if str(name).strip().lower() == "all responses"),
        None,
    )

    if target_sheet_name is not None:
        raw_df = workbook[target_sheet_name]
    else:
        raw_df = max(workbook.values(), key=lambda frame: frame.shape[0] * frame.shape[1])

    if raw_df.shape[0] < 4:
        raise ValueError("KAI rawdata 형식이 아닙니다. 최소 4행이 필요합니다.")

    question_row = raw_df.iloc[1].fillna("").astype(str).str.strip()
    type_row = raw_df.iloc[2].fillna("").astype(str).str.strip()

    data_df = raw_df.iloc[3:].copy().reset_index(drop=True)
    data_df = data_df.dropna(how="all")
    data_df.columns = list(range(raw_df.shape[1]))

    question_row.index = data_df.columns
    type_row.index = data_df.columns

    score_cols = [col for col in data_df.columns if "점수" in str(type_row[col])]
    if len(score_cols) < 10:
        score_like_cols: List[int] = []
        for col in data_df.columns:
            series = pd.to_numeric(data_df[col], errors="coerce").dropna()
            if not series.empty and series.between(1, 5).all():
                score_like_cols.append(col)
        score_cols = score_like_cols

    if len(score_cols) < 10:
        raise ValueError("KAI rawdata 형식이 아닙니다. 5점 척도 점수 컬럼 10개를 찾지 못했습니다.")

    score_cols = score_cols[:10]

    validity_col = next((col for col in data_df.columns if "적격성" in str(question_row[col])), None)
    if validity_col is not None:
        validity_values = data_df[validity_col].fillna("").astype(str).str.strip()
        valid_mask = validity_values.eq("적격")
        if valid_mask.any():
            data_df = data_df.loc[valid_mask].copy()

    module_col = next(
        (col for col in data_df.columns if "도움" in str(question_row[col]) and "무엇" in str(question_row[col])),
        None,
    )
    if module_col is None:
        fallback_col = score_cols[5] - 1
        if fallback_col in data_df.columns:
            module_col = fallback_col

    qualitative_cols: List[int] = []

    def find_question_col(*keywords: str) -> Optional[int]:
        for col in data_df.columns:
            text = str(question_row[col])
            if all(keyword in text for keyword in keywords):
                return col
        return None

    for keyset in [("긍정적으로", "느낀"), ("새롭게", "알게"), ("개선",)]:
        found_col = find_question_col(*keyset)
        if found_col is not None and found_col not in qualitative_cols:
            qualitative_cols.append(found_col)

    if len(qualitative_cols) < 3:
        for col in data_df.columns:
            if col <= score_cols[-1]:
                continue
            question_text = str(question_row[col]).strip()
            if not question_text or "적격" in question_text:
                continue

            series = data_df[col].fillna("").astype(str).str.strip()
            non_empty = series[series != ""]
            if non_empty.empty:
                continue

            numeric_ratio = pd.to_numeric(non_empty, errors="coerce").notna().mean()
            if numeric_ratio <= 0.2 and col not in qualitative_cols:
                qualitative_cols.append(col)
            if len(qualitative_cols) == 3:
                break

    if len(qualitative_cols) < 3:
        raise ValueError("KAI rawdata 형식이 아닙니다. 정성응답 컬럼 3개를 찾지 못했습니다.")

    return data_df, score_cols, module_col, qualitative_cols[:3]


def split_multi_select(text: str) -> List[str]:
    """Split a (복수 선택 가능) answer into individual selections.

    Commas inside brackets are kept intact so module names with internal commas
    (예: "기획의 표현 (Writing, Editing 등)") aren't broken apart. Newlines and
    semicolons are also treated as separators.
    """
    parts: List[str] = []
    buffer: List[str] = []
    depth = 0
    for ch in text:
        if ch in "([{":
            depth += 1
            buffer.append(ch)
        elif ch in ")]}":
            depth = max(0, depth - 1)
            buffer.append(ch)
        elif depth == 0 and ch in ",;\n":
            parts.append("".join(buffer).strip())
            buffer = []
        else:
            buffer.append(ch)
    parts.append("".join(buffer).strip())
    return [part for part in parts if part]


def summarize_module_preferences(series: pd.Series) -> List[Tuple[str, int]]:
    grouped: Dict[str, Dict[str, object]] = {}

    for raw in series.tolist():
        if pd.isna(raw):
            continue

        # 복수응답은 콤마로 구분되므로 개별 모듈로 분리해 각각 1명씩 집계한다.
        for piece in split_multi_select(str(raw)):
            text = re.sub(r"\s+", " ", piece).strip()
            if not text:
                continue

            key = re.sub(r"[^0-9a-z가-힣]+", "", text.lower())
            if not key:
                continue

            bucket = grouped.setdefault(key, {"count": 0, "labels": {}})
            bucket["count"] = int(bucket["count"]) + 1
            labels = bucket["labels"]
            labels[text] = labels.get(text, 0) + 1

    summarized: List[Tuple[str, int]] = []
    for bucket in grouped.values():
        labels = bucket["labels"]
        label = sorted(labels.items(), key=lambda item: (item[1], len(item[0])), reverse=True)[0][0]
        summarized.append((label, int(bucket["count"])))

    summarized.sort(key=lambda item: (item[1], len(item[0])), reverse=True)
    return summarized


def append_table_row_preserve_style(table, source_row_idx: int = -1) -> None:
    tbl = table._tbl
    source_tr = tbl.tr_lst[source_row_idx]
    new_tr = deepcopy(source_tr)

    for node in new_tr.iter():
        if str(node.tag).endswith("}t"):
            node.text = ""

    tbl.append(new_tr)


def update_kai_module_table(shape, module_rankings: List[Tuple[str, int]]) -> None:
    table = shape.table
    required_rows = 1 + len(module_rankings)

    while len(table.rows) < required_rows:
        append_table_row_preserve_style(table, -1)

    # 동점은 같은 순위로 표기한다 (예: 1, 2, 2, 4위).
    prev_count: Optional[int] = None
    prev_rank = 0
    for position, (name, count) in enumerate(module_rankings, start=1):
        if count == prev_count:
            rank = prev_rank
        else:
            rank = position
            prev_rank = rank
            prev_count = count
        if len(table.columns) > 0:
            set_text_preserve_style(table.cell(position, 0).text_frame, f"{rank}위")
        if len(table.columns) > 1:
            set_text_preserve_style(table.cell(position, 1).text_frame, f"{name}({count}명)")

    for row_idx in range(len(module_rankings) + 1, len(table.rows)):
        if len(table.columns) > 0:
            set_text_preserve_style(table.cell(row_idx, 0).text_frame, "")
        if len(table.columns) > 1:
            set_text_preserve_style(table.cell(row_idx, 1).text_frame, "")


def update_kai_qualitative_table(shape, df: pd.DataFrame, qualitative_cols: List[int]) -> None:
    table = shape.table
    for idx, col in enumerate(qualitative_cols):
        if len(table.rows) > 1 and len(table.columns) > idx:
            summary_text = summarize_qualitative_responses(df[col])
            set_text_preserve_style(table.cell(1, idx).text_frame, summary_text)


def update_kai_chart(shape, chart_idx: int, question_avgs: List[float]) -> None:
    def average(values: List[float]) -> float:
        valid = [value for value in values if pd.notna(value)]
        return float(sum(valid) / len(valid)) if valid else 0.0

    chart_categories = {
        0: ["과정 만족도", "운영 만족도", "강사 만족도"],
        1: ["전반 만족도", "동료 추천도", "기대 충족도", "학습목표 달성도", "현업 적용도"],
        2: ["운영자 만족도", "시설 만족도", "식사 만족도"],
        3: ["강의내용 만족도", "강의방식 만족도"],
    }

    chart_values = {
        0: [average(question_avgs[0:5]), average(question_avgs[5:8]), average(question_avgs[8:10])],
        1: question_avgs[0:5],
        2: question_avgs[5:8],
        3: question_avgs[8:10],
    }

    categories = chart_categories.get(chart_idx)
    values = chart_values.get(chart_idx)
    if categories is None or values is None:
        return

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("계열 1", tuple(round(float(value), 2) for value in values))
    shape.chart.replace_data(chart_data)
    set_chart_number_format(shape.chart, "0.00")


def classify_level(score: float, peers: Optional[List[float]] = None) -> str:
    rounded = round(float(score), 1)
    if rounded >= 4.6:
        base = "높은 수준"
    elif rounded == 4.5:
        base = "무난한 수준"
    elif rounded < 4.4:
        base = "비교적 낮은 수준"
    else:
        base = "무난한 수준"

    if peers:
        peer_max = max(peers)
        if rounded >= 4.4 and (peer_max - rounded) >= 0.1:
            return "비교적 낮은 수준"

    return base

def format_score_range(scores: List[float]) -> str:
    rounded = sorted(round(float(score), 1) for score in scores)
    if not rounded:
        return "0.0점"
    if rounded[0] == rounded[-1]:
        return f"{rounded[0]:.1f}점"
    return f"{rounded[0]:.1f}~{rounded[-1]:.1f}점"

def _has_batchim(word: str) -> bool:
    """Whether a Korean word ends in a final consonant (받침)."""
    if not word:
        return False
    last = word[-1]
    if "가" <= last <= "힣":
        return (ord(last) - 0xAC00) % 28 != 0
    return False


def _join_with_wa(names: List[str]) -> str:
    """Join Korean labels: 1 → as-is, 2 → 'A와/과 B', 3+ → comma list."""
    names = [name for name in names if name]
    if not names:
        return ""
    if len(names) == 1:
        return names[0]
    if len(names) == 2:
        particle = "과" if _has_batchim(names[0]) else "와"
        return f"{names[0]}{particle} {names[1]}"
    return ", ".join(names)


# 설문 문항 라벨 (5점 척도 10문항: q1~q5 과정, q6~q8 운영, q9~q10 강사)
COURSE_LABELS = ["전반 만족도", "동료 추천도", "기대 충족도", "학습목표 달성도", "현업 적용도"]
OPERATION_LABELS = ["운영자 만족도", "시설 만족도", "식사 만족도"]
INSTRUCTOR_LABELS = ["강의내용 만족도", "강의방식 만족도"]


def _describe_score_group(items: List[Tuple[str, float]]) -> str:
    """Name the highest- and lowest-scoring items in a section, chosen from
    actual data, with their level and score range."""
    scored = [(label, float(score)) for label, score in items]
    rounded = {label: round(score, 1) for label, score in scored}
    max_r = max(rounded.values())
    min_r = min(rounded.values())

    if max_r == min_r:
        level = classify_level(sum(score for _, score in scored) / len(scored))
        labels = [label for label, _ in scored]
        return f"{_join_with_wa(labels)}가 모두 {level}({format_score_range([s for _, s in scored])})으로 고르게 나타남."

    high = [label for label, _ in scored if rounded[label] == max_r]
    low = [label for label, _ in scored if rounded[label] == min_r]
    high_scores = [score for label, score in scored if rounded[label] == max_r]
    low_scores = [score for label, score in scored if rounded[label] == min_r]
    high_level = classify_level(max(high_scores))

    return (
        f"{_join_with_wa(high)}는 {high_level}({format_score_range(high_scores)})이었으며, "
        f"{_join_with_wa(low)}는 상대적으로 낮은 수준({format_score_range(low_scores)})으로 나타남."
    )


def build_kai_textbox4_paragraphs(
    question_avgs: List[float], module_rankings: List[Tuple[str, int]]
) -> List[str]:
    """Build all 7 paragraphs of slide 7 'TextBox 4' (시사점) from actual data:
    a headline + a data sentence for each of 과정/운영/강사, plus a module line."""
    scores = [float(value) for value in question_avgs]
    course = list(zip(COURSE_LABELS, scores[0:5]))
    operation = list(zip(OPERATION_LABELS, scores[5:8]))
    instructor = list(zip(INSTRUCTOR_LABELS, scores[8:10]))

    def avg(pairs: List[Tuple[str, float]]) -> float:
        return sum(score for _, score in pairs) / len(pairs) if pairs else 0.0

    def lowest_label(pairs: List[Tuple[str, float]]) -> str:
        return min(pairs, key=lambda item: item[1])[0]

    # 과정 만족도
    course_level = classify_level(avg(course))
    p0 = (
        f"과정 만족도 : 전반적인 과정 만족도가 {course_level}으로 나타났으며, "
        f"{lowest_label(course)}에서 소폭 개선 여지가 확인됨."
    )
    p1 = _describe_score_group(course)

    # 모듈 선호 (복수 선택 가능)
    if module_rankings:
        top_name, top_count = module_rankings[0]
        p2 = (
            f"'{top_name}' 모듈이 {top_count}명의 선택으로 선호도 1위를 기록했으며, "
            f"복수 선택 응답 기준 학습 활용도가 높은 것으로 확인됨."
        )
    else:
        p2 = "모듈별 선호 응답이 충분하지 않아 별도 분석은 생략함."

    # 운영 만족도
    operation_level = classify_level(avg(operation))
    p3 = (
        f"운영 만족도 : 운영 전반은 {operation_level}으로 평가되었으며, "
        f"{lowest_label(operation)} 등에서 개선 여지가 확인됨."
    )
    p4 = _describe_score_group(operation)

    # 강사 만족도
    instructor_level = classify_level(avg(instructor))
    p5 = f"강사 만족도 : 강사의 전문성과 강의 방식에 대한 학습자 수용성이 {instructor_level}으로 나타남."
    q9, q10 = scores[8], scores[9]
    p6 = (
        f"강의 내용({q9:.1f}점)과 강의방식({q10:.1f}점)의 만족도를 검토할 때 "
        f"교육 구성과 진행 방식에 대한 학습자들의 수용성이 {instructor_level}임을 확인할 수 있음."
    )

    return [p0, p1, p2, p3, p4, p5, p6]


def _set_paragraph_text(paragraph, value: str) -> None:
    """Replace a paragraph's text while keeping its first run's formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = value
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = value


def _find_slide7_textbox4(prs: Presentation):
    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide_idx != 7:
            continue
        fallback = None
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            shape_name = (shape.name or "").strip().lower()
            current_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if fallback is None and (
                "과정 만족도" in current_text and "운영 만족도" in current_text and "강사 만족도" in current_text
            ):
                fallback = shape
            if shape_name in {"textbox 4", "text box 4"}:
                return shape
        return fallback
    return None


def update_kai_slide7_textbox4(prs: Presentation, paragraphs: List[str]) -> None:
    target = _find_slide7_textbox4(prs)
    if target is None:
        return

    text_frame = target.text_frame
    while len(text_frame.paragraphs) < len(paragraphs):
        text_frame.add_paragraph()

    for idx, paragraph in enumerate(text_frame.paragraphs):
        value = paragraphs[idx] if idx < len(paragraphs) else ""
        _set_paragraph_text(paragraph, value)


def update_kai_summary_average(prs: Presentation, overall_avg: float) -> None:
    """Fill slide 3 '평균 : NN' textbox with the overall satisfaction average."""
    pattern = re.compile(r"(평균\s*[:：]\s*)\d+(?:\.\d+)?")
    replacement = f"평균 : {overall_avg:.2f}"
    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide_idx != 3:
            continue
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if pattern.search(paragraph.text):
                    _set_paragraph_text(paragraph, pattern.sub(replacement, paragraph.text))


def _fill_attendance_text(text: str, total: int, attended: int, rate: int) -> str:
    """Fill the '... 00명 中 ... 00명 참석 ... 00%' placeholders in one text block.

    Handles the wording variants used across slides (총 인원 / 전체 대상, 참석률 : / no
    label). Replaces the attended count first (the '00명' right before '참석'), then
    the remaining total '00명', then the rate '00%'. Each placeholder is two-or-more
    zeros, so real numbers are never touched.
    """
    if "참석" not in text or "%" not in text or "명" not in text:
        return text
    filled = re.sub(r"0{2,}(\s*)명(?=\s*참석)", rf"{attended}\1명", text, count=1)
    filled = re.sub(r"0{2,}(\s*)명", rf"{total}\1명", filled, count=1)
    filled = re.sub(r"0{2,}(\s*)%", rf"{rate}\1%", filled, count=1)
    return filled


def fill_attendance(prs: Presentation, total_count: Optional[int], attended_count: Optional[int]) -> None:
    """Auto-fill the repeated '총 인원 00명 中 00명 참석 (참석률 00%)' phrase.

    When 참석 인원 is left blank, full attendance (100%) is assumed.
    """
    try:
        total = int(total_count) if total_count else 0
    except (TypeError, ValueError):
        total = 0
    if total <= 0:
        return

    try:
        attended = int(attended_count) if attended_count else 0
    except (TypeError, ValueError):
        attended = 0
    if attended <= 0:
        attended = total
    attended = min(attended, total)
    rate = int(round((attended / total) * 100))

    def process(text_frame) -> None:
        for paragraph in text_frame.paragraphs:
            new_text = _fill_attendance_text(paragraph.text, total, attended, rate)
            if new_text != paragraph.text:
                _set_paragraph_text(paragraph, new_text)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                process(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        process(cell.text_frame)


def populate_ppt_lam(excel_bytes: bytes, class_name: str, template_path: Path,
                     total_count: Optional[int] = None, attended_count: Optional[int] = None) -> bytes:
    df = pd.read_excel(io.BytesIO(excel_bytes))
    question_cols = identify_question_columns(df)

    placeholders, percentages_by_question, counts_by_question, respondent_count = compute_lam_metrics(df, question_cols)

    replacements: Dict[str, str] = {key: format_score(value) for key, value in placeholders.items()}
    replacements["class_name"] = class_name.strip() if class_name.strip() else "과정명 미입력"
    replacements["respondent_count"] = str(respondent_count)

    prs = Presentation(str(template_path))

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            name = (shape.name or "").strip().lower()

            if shape.has_chart:
                chart_idx = shape_index_from_name(name)
                if chart_idx == 0:
                    update_lam_chart_0(shape, placeholders)
                elif chart_idx is not None and 1 <= chart_idx <= 10:
                    update_lam_question_chart(shape, chart_idx, percentages_by_question)

            if shape.has_table:
                table_idx = shape_index_from_name(name)
                if table_idx is None:
                    continue
                if 1 <= table_idx <= 10:
                    update_lam_question_table(shape, table_idx, counts_by_question)
                    format_table_font(shape.table, "Noto Sans CJK KR DemiLight", 9)
                elif table_idx == 11:
                    update_lam_qualitative_table(shape, df)
                    format_table_font(shape.table, "Noto Sans CJK KR DemiLight", 9)

    replace_text_placeholders(prs, replacements)
    fill_attendance(prs, total_count, attended_count)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def populate_ppt_kai(excel_bytes: bytes, class_name: str, template_path: Path,
                     total_count: Optional[int] = None, attended_count: Optional[int] = None) -> bytes:
    df, score_cols, module_col, qualitative_cols = parse_kai_excel(excel_bytes)

    numeric = df[score_cols].apply(pd.to_numeric, errors="coerce")
    respondent_count = int(numeric.dropna(how="all").shape[0])
    if respondent_count == 0:
        raise ValueError("유효 응답자가 없습니다.")

    question_avgs: List[float] = []
    for col in score_cols:
        value = numeric[col].mean()
        question_avgs.append(float(value) if pd.notna(value) else 0.0)

    module_rankings: List[Tuple[str, int]] = []
    if module_col is not None:
        module_rankings = summarize_module_preferences(df[module_col])

    textbox4_paragraphs = build_kai_textbox4_paragraphs(question_avgs, module_rankings)
    overall_avg = sum(question_avgs) / len(question_avgs) if question_avgs else 0.0

    replacements = {
        "class_name": class_name.strip() if class_name.strip() else "과정명 미입력",
        "n_number": str(respondent_count),
        "respondent_count": str(respondent_count),
    }
    prs = Presentation(str(template_path))

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            name = (shape.name or "").strip().lower()

            if shape.has_chart:
                chart_idx = shape_index_from_name(name)
                if chart_idx is not None and 0 <= chart_idx <= 3:
                    update_kai_chart(shape, chart_idx, question_avgs)

            if shape.has_table:
                table_idx = shape_index_from_name(name)
                if table_idx == 0:
                    update_kai_module_table(shape, module_rankings)
                elif slide_idx == 8 and table_idx == 6 and len(shape.table.columns) >= 3:
                    update_kai_qualitative_table(shape, df, qualitative_cols)
                    format_table_font(shape.table, "Noto Sans CJK KR DemiLight", 9)

    update_kai_slide7_textbox4(prs, textbox4_paragraphs)
    update_kai_summary_average(prs, overall_avg)
    replace_text_placeholders(prs, replacements)
    fill_attendance(prs, total_count, attended_count)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def populate_ppt(excel_bytes: bytes, class_name: str, report_type: str, template_path: Path,
                 total_count: Optional[int] = None, attended_count: Optional[int] = None) -> bytes:
    if report_type == "lam":
        return populate_ppt_lam(excel_bytes, class_name, template_path, total_count, attended_count)
    if report_type == "kai":
        return populate_ppt_kai(excel_bytes, class_name, template_path, total_count, attended_count)
    raise ValueError(f"지원하지 않는 보고서 유형입니다: {report_type}")


def main() -> None:
    st.set_page_config(page_title="결과보고서 PPT 자동 생성기", layout="centered")
    st.title("결과보고서 PPT 자동 생성기")
    st.write("raw data 엑셀을 업로드하면 템플릿 기반 결과보고서를 자동 생성합니다.")

    report_label = st.selectbox("보고서 유형", options=list(REPORT_TYPES.keys()), index=0)
    report_type = str(REPORT_TYPES[report_label]["id"])

    try:
        template_path = resolve_template_path(report_label)
        st.caption(f"사용 템플릿: {template_path}")
    except FileNotFoundError as exc:
        template_path = None
        st.error(str(exc))

    class_name = st.text_input("과정명", placeholder="예: 2026년 신입사원 교육")

    attendance_col1, attendance_col2 = st.columns(2)
    with attendance_col1:
        total_count = st.number_input("총 인원", min_value=0, value=0, step=1,
                                      help="참석자 명단·참석 현황 문구에 자동으로 채워집니다. (0이면 미입력)")
    with attendance_col2:
        attended_count = st.number_input("참석 인원", min_value=0, value=0, step=1,
                                         help="비워두면 총 인원과 동일(참석률 100%)로 계산됩니다.")

    uploaded_excel = st.file_uploader("원본(raw data) 파일 업로드 (.xlsx)", type=["xlsx"])

    if st.button("PPT 생성", type="primary"):
        if not uploaded_excel:
            st.error("원본(raw data) 파일을 먼저 업로드해 주세요.")
            return
        if template_path is None:
            st.error("선택한 보고서 유형의 템플릿 파일을 찾지 못했습니다.")
            return

        try:
            ppt_bytes = populate_ppt(
                uploaded_excel.read(),
                class_name,
                report_type,
                template_path,
                total_count=int(total_count) or None,
                attended_count=int(attended_count) or None,
            )
        except Exception as exc:  # noqa: BLE001
            st.exception(exc)
            return

        st.success("PPT가 생성되었습니다. 아래 버튼으로 다운로드해 주세요.")
        st.download_button(
            label="결과 PPT 다운로드",
            data=ppt_bytes,
            file_name=f"{class_name.strip() or report_type}_만족도_보고서.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )


if __name__ == "__main__":
    main()
